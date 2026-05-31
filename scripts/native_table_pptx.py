
#!/usr/bin/env python3
"""Build a one-slide PPTX containing a native PowerPoint table from JSON.

Spec example:
{
  "slide_size": [13.333, 7.5],
  "table_box": [0.5, 1.2, 6.0, 4.8],
  "headers": ["序号", "指标", "增长率（%）"],
  "rows": [["1", "地区生产总值 GDP", "3.5%"]],
  "font": "Microsoft YaHei",
  "header_fill": "082D5A",
  "header_font_color": "FFFFFF",
  "body_font_color": "111111"
}
All coordinates are inches.
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def rgb(value, default="FFFFFF"):
    value = (value or default).strip().lstrip("#")
    return RGBColor.from_string(value.upper())


def set_cell_text(cell, text, font_name, font_size, color, bold=False):
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def build(spec, out_path):
    prs = Presentation()
    slide_w, slide_h = spec.get("slide_size", [13.333, 7.5])
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    data = [headers] + rows if headers else rows
    if not data or not data[0]:
        raise SystemExit("spec must contain headers or rows")

    x, y, w, h = spec.get("table_box", [0.5, 0.5, slide_w - 1, slide_h - 1])
    table_shape = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    font = spec.get("font", "Microsoft YaHei")
    body_size = float(spec.get("body_font_size", 12))
    header_size = float(spec.get("header_font_size", body_size + 1))

    header_fill = rgb(spec.get("header_fill", "082D5A"))
    header_color = rgb(spec.get("header_font_color", "FFFFFF"))
    body_color = rgb(spec.get("body_font_color", "111111"))

    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            if r == 0 and headers:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                set_cell_text(cell, value, font, header_size, header_color, True)
            else:
                set_cell_text(cell, value, font, body_size, body_color, False)

    col_widths = spec.get("col_widths")
    if col_widths:
        for idx, width in enumerate(col_widths[: len(data[0])]):
            table.columns[idx].width = Inches(float(width))

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", required=True)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()
    build(json.loads(Path(args.spec).read_text(encoding="utf-8")), args.out)


if __name__ == "__main__":
    main()
